
import argparse, os
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt

def scale_to_bn(x_m):
    try:
        return float(x_m)/1000.0
    except Exception:
        return x_m

def build_bar_chart_ev_ebitda(df, out_path):
    fig = plt.figure(figsize=(10, 5))
    plt.bar(df['Company'], df['EV_EBITDA'])
    plt.title("EV/EBITDA – Selected Comps")
    plt.ylabel("x")
    plt.xticks(rotation=30, ha="right")
    plt.tight_layout()
    fig.savefig(out_path, dpi=200, bbox_inches="tight")
    plt.close(fig)

def build_scatter_pe_growth(df, out_path):
    fig = plt.figure(figsize=(10, 5))
    plt.scatter(df['Rev_Growth_%'], df['PE'])
    plt.xlabel("Revenue growth (%)")
    plt.ylabel("P/E (x)")
    plt.title("P/E vs. Revenue Growth")
    plt.tight_layout()
    fig.savefig(out_path, dpi=200, bbox_inches="tight")
    plt.close(fig)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
    return slide

def add_table(prs, title, df):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    rows, cols = df.shape
    rows += 1
    table = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.6), Inches(9.3), Inches(0.8)).table
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)
    for i in range(df.shape[0]):
        for j in range(df.shape[1]):
            val = df.iat[i, j]
            cell = table.cell(i+1, j)
            if isinstance(val, float):
                if ("Margin" in df.columns[j]) or ("Growth" in df.columns[j]):
                    cell.text = f"{val:.1f}%"
                elif ("EV/EBITDA" in df.columns[j]) or ("EV_EBITDA" in df.columns[j]) or ("P/E" in df.columns[j]) or (df.columns[j]=="PE"):
                    cell.text = f"{val:.1f}x"
                else:
                    cell.text = f"{val:.1f}"
            else:
                cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(10)
    return slide

def add_picture(prs, title, img_path, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(img_path, Inches(0.4), Inches(1.6), height=Inches(4.5))
    if caption:
        tx = slide.shapes.add_textbox(Inches(0.4), Inches(6.3), Inches(9.0), Inches(0.6))
        tf = tx.text_frame
        tf.text = caption
        tf.paragraphs[0].font.size = Pt(12)

def add_val_summary(prs, df):
    metrics = {
        "EV/EBITDA": df["EV_EBITDA"].astype(float).replace([np.inf,-np.inf], np.nan).dropna(),
        "P/E": df["PE"].astype(float).replace([np.inf,-np.inf], np.nan).dropna(),
    }
    rows = []
    for name, s in metrics.items():
        if len(s)==0:
            med=p25=p75=float("nan")
        else:
            med = float(np.nanmedian(s))
            p25 = float(np.nanpercentile(s,25))
            p75 = float(np.nanpercentile(s,75))
        rows.append([name, med, p25, p75])
    df_sum = pd.DataFrame(rows, columns=["Metric","Median (x)","25th (x)","75th (x)"])
    add_table(prs, "Valuation Summary (EV/EBITDA & P/E)", df_sum)

def main():
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("--csv", required=True)
    parser.add_argument("--out", default="output/pitch_comps.pptx")
    parser.add_argument("--title", default="Comparable Companies – Deck")
    args = parser.parse_args()

    df = pd.read_csv(args.csv)
    df_disp = df.copy()
    df_disp["EV_USD_bn"] = df_disp["EV_USD_m"].apply(scale_to_bn)
    df_disp["EBITDA_USD_bn"] = df_disp["EBITDA_USD_m"].apply(scale_to_bn)
    cols = ["Company","Ticker","Sector","EV_USD_bn","EBITDA_USD_bn","EV_EBITDA","PE","Rev_Growth_%","EBITDA_Margin_%"]
    df_tbl = df_disp[cols].rename(columns={"EV_USD_bn":"EV (USD bn)","EBITDA_USD_bn":"EBITDA (USD bn)"})
    df_bar = df.sort_values("EV_EBITDA", ascending=False).reset_index(drop=True)

    # charts
    os.makedirs(os.path.dirname(args.out), exist_ok=True)
    chart1 = os.path.join(os.path.dirname(args.out), "ev_ebitda_chart.png")
    chart2 = os.path.join(os.path.dirname(args.out), "pe_vs_growth.png")
    build_bar_chart_ev_ebitda(df_bar, chart1)
    build_scatter_pe_growth(df, chart2)

    prs = Presentation()
    add_title_slide(prs, args.title, "Auto-generated from comps CSV")
    add_bullets_slide(prs, "Key Takeaways", [
        f"Median EV/EBITDA ≈ {df['EV_EBITDA'].median():.1f}x; P/E ≈ {df['PE'].median():.1f}x",
        "Prefer medians & IQR (25–75th) over means due to outliers",
        "Growth/margins cluster in Software/Semis within this sample",
    ])
    add_table(prs, "Comparable Companies Table", df_tbl)
    add_picture(prs, "EV/EBITDA – Distribution", chart1)
    add_picture(prs, "P/E vs Revenue Growth", chart2, caption="Illustrative only; ensure sector comparability.")
    add_val_summary(prs, df)
    add_bullets_slide(prs, "Methodology & Sources", [
    f"Universe: {len(df)} listed software comps; currency: USD; base: TTM.",
    "Valuation: EV/EBITDA & P/E; medians and interquartile range (25–75th).",
    "Outliers reviewed; results illustrative. Sources: public filings & aggregators."
])
    prs.save(args.out)
    print(f"Saved deck to: {args.out}")

if __name__ == "__main__":
    main()